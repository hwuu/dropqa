"""文档解析器测试"""

import uuid
from pathlib import Path
import tempfile

import pytest

from dropqa.indexer.parser import (
    MarkdownParser,
    ParsedNode,
    flatten_nodes,
    parse_document,
    SUPPORTED_EXTENSIONS,
    _filter_empty_sections,
)


class TestMarkdownParser:
    """Markdown 解析器测试"""

    def test_parse_empty_document(self):
        """测试解析空文档"""
        parser = MarkdownParser()
        result = parser.parse("", "test.md")

        assert result.node_type == "document"
        assert result.depth == 0
        assert result.title == "test.md"
        assert result.children == []

    def test_parse_simple_heading(self):
        """测试解析简单标题"""
        parser = MarkdownParser()
        content = "# Hello World"
        result = parser.parse(content, "test.md")

        assert len(result.children) == 1
        section = result.children[0]
        assert section.node_type == "section"
        assert section.depth == 1
        assert section.title == "Hello World"

    def test_parse_multiple_headings(self):
        """测试解析多个同级标题"""
        parser = MarkdownParser()
        content = """# First
# Second
# Third"""
        result = parser.parse(content, "test.md")

        assert len(result.children) == 3
        assert result.children[0].title == "First"
        assert result.children[1].title == "Second"
        assert result.children[2].title == "Third"

    def test_parse_nested_headings(self):
        """测试解析嵌套标题"""
        parser = MarkdownParser()
        content = """# Chapter 1
## Section 1.1
### Subsection 1.1.1
## Section 1.2
# Chapter 2"""
        result = parser.parse(content, "test.md")

        # 根节点有 2 个 chapter
        assert len(result.children) == 2

        chapter1 = result.children[0]
        assert chapter1.title == "Chapter 1"
        assert len(chapter1.children) == 2  # Section 1.1 和 Section 1.2

        section11 = chapter1.children[0]
        assert section11.title == "Section 1.1"
        assert len(section11.children) == 1  # Subsection 1.1.1

        subsection = section11.children[0]
        assert subsection.title == "Subsection 1.1.1"

    def test_parse_with_paragraphs(self):
        """测试解析带段落的文档"""
        parser = MarkdownParser()
        content = """# Introduction

This is the introduction paragraph.

## Details

Here are the details.
Multiple lines in the same paragraph."""
        result = parser.parse(content, "test.md")

        # 检查结构
        assert len(result.children) == 1  # Introduction

        intro = result.children[0]
        assert intro.title == "Introduction"
        assert len(intro.children) == 2  # paragraph + Details

        # 第一个子节点是段落
        para1 = intro.children[0]
        assert para1.node_type == "paragraph"
        assert "introduction paragraph" in para1.content

        # 第二个子节点是 Details section
        details = intro.children[1]
        assert details.node_type == "section"
        assert details.title == "Details"
        assert len(details.children) == 1  # paragraph

    def test_parse_content_before_first_heading(self):
        """测试第一个标题前的内容"""
        parser = MarkdownParser()
        content = """Some content before heading.

# First Heading"""
        result = parser.parse(content, "test.md")

        # 第一个子节点是段落（挂在根节点下）
        assert len(result.children) == 2
        assert result.children[0].node_type == "paragraph"
        assert result.children[1].node_type == "section"

    def test_heading_levels(self):
        """测试各级标题"""
        parser = MarkdownParser()
        content = """# H1
## H2
### H3
#### H4
##### H5
###### H6"""
        result = parser.parse(content, "test.md")

        # 遍历检查深度
        def check_depth(node, expected_depth):
            assert node.depth == expected_depth
            for child in node.children:
                if child.node_type == "section":
                    check_depth(child, expected_depth + 1)

        check_depth(result.children[0], 1)


class TestFlattenNodes:
    """节点展平测试"""

    def test_flatten_simple(self):
        """测试简单展平"""
        root = ParsedNode(
            node_type="document",
            depth=0,
            title="test.md",
        )
        root.children.append(
            ParsedNode(
                node_type="section",
                depth=1,
                title="Chapter 1",
            )
        )

        doc_id = uuid.uuid4()
        nodes = flatten_nodes(root, doc_id)

        assert len(nodes) == 2
        assert nodes[0]["node_type"] == "document"
        assert nodes[0]["parent_id"] is None
        assert nodes[1]["node_type"] == "section"
        assert nodes[1]["parent_id"] == nodes[0]["id"]

    def test_flatten_preserves_structure(self):
        """测试展平保留结构"""
        parser = MarkdownParser()
        content = """# Chapter 1
## Section 1.1
Paragraph content"""
        root = parser.parse(content, "test.md")

        doc_id = uuid.uuid4()
        nodes = flatten_nodes(root, doc_id)

        # document, section, section, paragraph
        assert len(nodes) == 4

        # 检查类型
        types = [n["node_type"] for n in nodes]
        assert types == ["document", "section", "section", "paragraph"]

        # 检查父子关系
        doc_node = nodes[0]
        chapter_node = nodes[1]
        section_node = nodes[2]
        para_node = nodes[3]

        assert chapter_node["parent_id"] == doc_node["id"]
        assert section_node["parent_id"] == chapter_node["id"]
        assert para_node["parent_id"] == section_node["id"]


class TestParseDocument:
    """parse_document 工厂函数测试"""

    def test_supported_extensions(self):
        """测试支持的扩展名"""
        assert ".md" in SUPPORTED_EXTENSIONS
        assert ".docx" in SUPPORTED_EXTENSIONS
        assert ".pptx" in SUPPORTED_EXTENSIONS
        assert ".xlsx" in SUPPORTED_EXTENSIONS

    def test_parse_markdown_via_factory(self, tmp_path):
        """测试通过工厂函数解析 Markdown"""
        md_file = tmp_path / "test.md"
        md_file.write_text("# Hello\n\nWorld", encoding="utf-8")

        result = parse_document(md_file)

        assert result.node_type == "document"
        assert result.title == "test.md"
        assert len(result.children) == 1
        assert result.children[0].title == "Hello"

    def test_unsupported_extension(self, tmp_path):
        """测试不支持的扩展名"""
        txt_file = tmp_path / "test.txt"
        txt_file.write_text("Hello", encoding="utf-8")

        with pytest.raises(ValueError, match="不支持的文件格式"):
            parse_document(txt_file)


class TestFilterEmptySections:
    """_filter_empty_sections 辅助函数测试"""

    def test_filter_empty_section(self):
        """测试过滤空 section"""
        root = ParsedNode(
            node_type="document",
            depth=0,
            title="test.md",
        )
        # 添加一个空 section（只有标题，无内容）
        empty_section = ParsedNode(
            node_type="section",
            depth=1,
            title="Empty Section",
        )
        root.children.append(empty_section)

        _filter_empty_sections(root)

        assert len(root.children) == 0

    def test_keep_section_with_content(self):
        """测试保留有内容的 section"""
        root = ParsedNode(
            node_type="document",
            depth=0,
            title="test.md",
        )
        # 添加一个有内容的 section
        section = ParsedNode(
            node_type="section",
            depth=1,
            title="Section with Content",
        )
        para = ParsedNode(
            node_type="paragraph",
            depth=2,
            content="Some content",
        )
        section.children.append(para)
        root.children.append(section)

        _filter_empty_sections(root)

        assert len(root.children) == 1
        assert root.children[0].title == "Section with Content"

    def test_filter_nested_empty_sections(self):
        """测试过滤嵌套的空 section"""
        root = ParsedNode(
            node_type="document",
            depth=0,
            title="test.md",
        )
        # 创建嵌套的空 section
        section1 = ParsedNode(
            node_type="section",
            depth=1,
            title="Section 1",
        )
        section2 = ParsedNode(
            node_type="section",
            depth=2,
            title="Section 2",
        )
        section1.children.append(section2)
        root.children.append(section1)

        _filter_empty_sections(root)

        # 两个 section 都是空的，应该都被删除
        assert len(root.children) == 0

    def test_filter_empty_paragraph(self):
        """测试过滤空段落"""
        root = ParsedNode(
            node_type="document",
            depth=0,
            title="test.md",
        )
        section = ParsedNode(
            node_type="section",
            depth=1,
            title="Section",
        )
        empty_para = ParsedNode(
            node_type="paragraph",
            depth=2,
            content="",
        )
        valid_para = ParsedNode(
            node_type="paragraph",
            depth=2,
            content="Valid content",
        )
        section.children.append(empty_para)
        section.children.append(valid_para)
        root.children.append(section)

        _filter_empty_sections(root)

        # section 应该保留（因为有有效段落）
        assert len(root.children) == 1
        # 空段落应该被删除
        assert len(root.children[0].children) == 1
        assert root.children[0].children[0].content == "Valid content"


class TestParseDocx:
    """DOCX 解析测试（需要 unstructured 库）"""

    @pytest.fixture
    def sample_docx(self, tmp_path):
        """创建测试用 DOCX 文件"""
        try:
            from docx import Document
        except ImportError:
            pytest.skip("python-docx not installed")

        doc = Document()
        doc.add_heading("Chapter 1", level=1)
        doc.add_paragraph("This is the first paragraph.")
        doc.add_heading("Section 1.1", level=2)
        doc.add_paragraph("This is content under section 1.1")
        doc.add_heading("Chapter 2", level=1)
        doc.add_paragraph("Content for chapter 2")

        file_path = tmp_path / "test.docx"
        doc.save(str(file_path))
        return file_path

    def test_parse_docx_basic(self, sample_docx):
        """测试基本 DOCX 解析"""
        try:
            result = parse_document(sample_docx)
        except ImportError:
            pytest.skip("unstructured not installed")

        assert result.node_type == "document"
        assert result.title == "test.docx"
        # 检查解析出了内容
        assert len(result.children) > 0


class TestParsePptx:
    """PPTX 解析测试（需要 unstructured 库）"""

    @pytest.fixture
    def sample_pptx(self, tmp_path):
        """创建测试用 PPTX 文件"""
        try:
            from pptx import Presentation
            from pptx.util import Inches
        except ImportError:
            pytest.skip("python-pptx not installed")

        prs = Presentation()

        # 添加第一张幻灯片
        slide_layout = prs.slide_layouts[1]  # 标题和内容
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "Slide 1 Title"
        body = slide.shapes.placeholders[1]
        body.text = "Content for slide 1"

        # 添加第二张幻灯片
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "Slide 2 Title"
        body = slide.shapes.placeholders[1]
        body.text = "Content for slide 2"

        file_path = tmp_path / "test.pptx"
        prs.save(str(file_path))
        return file_path

    def test_parse_pptx_basic(self, sample_pptx):
        """测试基本 PPTX 解析"""
        try:
            result = parse_document(sample_pptx)
        except ImportError:
            pytest.skip("unstructured not installed")

        assert result.node_type == "document"
        assert result.title == "test.pptx"
        # 检查解析出了内容
        assert len(result.children) > 0


class TestParseXlsx:
    """XLSX 解析测试"""

    @pytest.fixture
    def sample_xlsx(self, tmp_path):
        """创建测试用 XLSX 文件"""
        try:
            from openpyxl import Workbook
        except ImportError:
            pytest.skip("openpyxl not installed")

        wb = Workbook()
        ws = wb.active
        ws.title = "Sheet1"

        # 添加表头
        ws["A1"] = "Name"
        ws["B1"] = "Age"
        ws["C1"] = "City"

        # 添加数据行
        ws["A2"] = "Alice"
        ws["B2"] = 25
        ws["C2"] = "Beijing"

        ws["A3"] = "Bob"
        ws["B3"] = 30
        ws["C3"] = "Shanghai"

        # 空行（第4行）
        # 再添加一行数据
        ws["A5"] = "Charlie"
        ws["B5"] = 35
        ws["C5"] = "Guangzhou"

        file_path = tmp_path / "test.xlsx"
        wb.save(str(file_path))
        return file_path

    def test_parse_xlsx_basic(self, sample_xlsx):
        """测试基本 XLSX 解析"""
        result = parse_document(sample_xlsx)

        assert result.node_type == "document"
        assert result.title == "test.xlsx"
        assert len(result.children) == 1  # 一个 sheet

        sheet = result.children[0]
        assert sheet.node_type == "section"
        assert sheet.title == "Sheet1"
        # 3 行数据（跳过表头和空行）
        assert len(sheet.children) == 3

    def test_parse_xlsx_content_format(self, sample_xlsx):
        """测试 XLSX 内容格式"""
        result = parse_document(sample_xlsx)

        sheet = result.children[0]
        first_row = sheet.children[0]

        assert first_row.node_type == "section"
        assert first_row.title == "第2行"
        # 检查内容格式：字段名: 值
        assert "Name: Alice" in first_row.content
        assert "Age: 25" in first_row.content
        assert "City: Beijing" in first_row.content

    def test_parse_xlsx_skip_empty_rows(self, sample_xlsx):
        """测试跳过空行"""
        result = parse_document(sample_xlsx)

        sheet = result.children[0]
        # 第4行是空行，应该被跳过
        row_titles = [child.title for child in sheet.children]
        assert "第4行" not in row_titles

    @pytest.fixture
    def multi_sheet_xlsx(self, tmp_path):
        """创建多 sheet 的 XLSX 文件"""
        try:
            from openpyxl import Workbook
        except ImportError:
            pytest.skip("openpyxl not installed")

        wb = Workbook()
        ws1 = wb.active
        ws1.title = "Users"
        ws1["A1"] = "Name"
        ws1["A2"] = "Alice"

        ws2 = wb.create_sheet("Products")
        ws2["A1"] = "Product"
        ws2["A2"] = "Apple"

        # 空 sheet
        ws3 = wb.create_sheet("Empty")

        file_path = tmp_path / "multi_sheet.xlsx"
        wb.save(str(file_path))
        return file_path

    def test_parse_xlsx_multi_sheet(self, multi_sheet_xlsx):
        """测试多 sheet 解析"""
        result = parse_document(multi_sheet_xlsx)

        # 只有 2 个非空 sheet
        assert len(result.children) == 2

        sheet_titles = [child.title for child in result.children]
        assert "Users" in sheet_titles
        assert "Products" in sheet_titles
        assert "Empty" not in sheet_titles  # 空 sheet 被跳过
